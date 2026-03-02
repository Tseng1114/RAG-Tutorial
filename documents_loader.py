import os
import glob

class pdf_loader:
    def load(self, file_path: str) -> str:
        import fitz
        doc = fitz.open(file_path)
        text = "".join([page.get_text() for page in doc])
        doc.close()
        return text


class docx_loader:
    def load(self, file_path: str) -> str:
        from docx import Document
        doc = Document(file_path)
        paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
        for table in doc.tables:
            for row in table.rows:
                row_text = " | ".join(cell.text.strip() for cell in row.cells if cell.text.strip())
                if row_text:
                    paragraphs.append(row_text)
        return "\n".join(paragraphs)


class pptx_loader:
    def load(self, file_path: str) -> str:
        from pptx import Presentation
        prs = Presentation(file_path)
        lines = []
        for slide_num, slide in enumerate(prs.slides, start=1):
            lines.append(f"[Slide {slide_num}]")
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        text = para.text.strip()
                        if text:
                            lines.append(text)
        return "\n".join(lines)


class excel_loader:
    def load(self, file_path: str) -> str:
        import openpyxl
        wb = openpyxl.load_workbook(file_path, data_only=True)
        lines = []
        for sheet in wb.worksheets:
            lines.append(f"[Sheet: {sheet.title}]")
            for row in sheet.iter_rows(values_only=True):
                row_text = " | ".join(str(cell) for cell in row if cell is not None)
                if row_text.strip():
                    lines.append(row_text)
        return "\n".join(lines)


class csv_loader:
    def load(self, file_path: str) -> str:
        import csv
        lines = []
        with open(file_path, newline='', encoding='utf-8-sig') as f:
            reader = csv.reader(f)
            for row in reader:
                row_text = " | ".join(cell.strip() for cell in row if cell.strip())
                if row_text:
                    lines.append(row_text)
        return "\n".join(lines)


class txt_loader:
    def load(self, file_path: str) -> str:
        with open(file_path, encoding='utf-8') as f:
            return f.read()

LOADER_MAP = {
    ".pdf":  pdf_loader(),
    ".docx": docx_loader(),
    ".pptx": pptx_loader(),
    ".xlsx": excel_loader(),
    ".xls":  excel_loader(),
    ".csv":  csv_loader(),
    ".txt":  txt_loader(),
}


def load_documents(directory: str) -> list[dict]:
    supported_extensions = set(LOADER_MAP.keys())
    all_files = []

    for ext in supported_extensions:
        pattern = os.path.join(directory, f"*{ext}")
        all_files.extend(glob.glob(pattern))

    if not all_files:
        print(f"No supported files found in: {directory}")
        print(f"Supported formats: {', '.join(sorted(supported_extensions))}")
        return []

    raw_docs = []

    for file_path in sorted(all_files):
        file_name = os.path.basename(file_path)
        ext = os.path.splitext(file_name)[1].lower()
        loader = LOADER_MAP.get(ext)

        print(f"Reading [{ext}]: {file_name}")

        try:
            content = loader.load(file_path)

            if content.strip():
                raw_docs.append({
                    "content": content,
                    "source": file_name
                })
            else:
                print(f"  Warning: No text extracted from {file_name}, skipping.")

        except Exception as e:
            print(f"  Error reading {file_name}: {e}")

    print(f"\nLoaded {len(raw_docs)} / {len(all_files)} document(s).")
    return raw_docs